import xml.etree.ElementTree as ET
import docx
import fitz  # PyMuPDF: PDF 텍스트 추출용
import magic
import zipfile
from pptx import Presentation
import olefile
import subprocess
import re
import os
from datetime import datetime, timedelta
import struct, zlib, pathlib
from pathlib import Path
import PyPDF2
import openpyxl
import xlrd


def to_iso_str(dt):
    """datetime 또는 문자열을 ISO 8601 문자열로 변환"""
    if isinstance(dt, datetime):
        return dt.strftime("%Y-%m-%dT%H:%M:%S")
    elif isinstance(dt, str):
        dt = dt.replace("D:", "").strip()
        try:
            parsed = datetime.strptime(dt[:14], "%Y%m%d%H%M%S")
            return parsed.strftime("%Y-%m-%dT%H:%M:%S")
        except:
            return ""
    return ""


def get_ole_created_date(path):
    try:
        if not olefile.isOleFile(path):
            return ""
        ole = olefile.OleFileIO(path)
        if ole.exists("\x05SummaryInformation"):
            props = ole.getproperties("\x05SummaryInformation")
            created = props.get(13)
            if isinstance(created, (int, float)):
                created = datetime(1601, 1, 1) + timedelta(seconds=created)
                return to_iso_str(created)
            else:
                return created
    except:
        return ""
    return ""


def get_office_openxml_created_date(path):
    try:
        with zipfile.ZipFile(path, "r") as z:
            core_xml = z.read("docProps/core.xml").decode("utf-8", "ignore")
            root = ET.fromstring(core_xml)
            for t_elem in root.iter("{http://purl.org/dc/terms/}modified"):
                if t_elem.text is not None:
                    return t_elem.text.replace("Z", "")
    except:
        return ""
    return ""


def get_hwpx_created_date(path):
    try:
        with zipfile.ZipFile(path, "r") as z:
            for name in z.namelist():
                if "FileHeader.xml" in name or "Contents/content.hpf" in name:
                    content = z.read(name).decode("utf-8", "ignore")
                    root = ET.fromstring(content)
                    for t_elem in root.iter("{http://www.idpf.org/2007/opf/}meta"):
                        if "ModifiedDate" in t_elem.attrib["name"]:
                            return t_elem.text.replace("Z", "")
    except:
        pass
    return ""


def get_pdf_created_date(path):
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            info = reader.metadata
            return to_iso_str(info.get("/CreationDate"))
    except:
        return ""


def get_file_ctime_created_date(path):
    try:
        ts = os.path.getctime(path)
        return to_iso_str(datetime.fromtimestamp(ts))
    except:
        return ""


def extract_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    chunks = []
    for ws in wb.worksheets:               # 모든 시트 순회
        for row in ws.iter_rows(values_only=True):
            for value in row:
                if value not in (None, ""):
                    chunks.append(str(value))
    wb.close()
    return "\n".join(chunks)


def extract_xls(filepath):
    book = xlrd.open_workbook(filepath, formatting_info=False)
    chunks = []
    for sheet in book.sheets():
        for r in range(sheet.nrows):
            for c in range(sheet.ncols):
                val = sheet.cell_value(r, c)
                if val not in ("", None):
                    chunks.append(str(val))
    return "\n".join(chunks)


def extract_ppt_doc(filepath):
    from org.apache.poi.sl.extractor import SlideShowExtractor
    from org.apache.poi.extractor import ExtractorFactory
    from java.io import FileInputStream

    try:
        with FileInputStream(filepath) as stream:
            extractor = ExtractorFactory.createExtractor(stream)  # 형식 자동 감지
            if isinstance(extractor, SlideShowExtractor):
                extractor.setSlidesByDefault(True)
                extractor.setNotesByDefault(True)
            ans = extractor.getText()[:None]
            ans = ans.replace("\n", " ")
        return ans
    except RuntimeError as e:
        return ""
    return ""


def extract_hwp(filepath):
    filepath = pathlib.Path(filepath).expanduser()
    with filepath.open("rb") as f:
        magic = f.read(4)
        f.seek(0)
        if magic == b"HWP ":
            header = f.read(256)
            body_off = struct.unpack_from("<I", header, 36)[0]
            f.seek(body_off)
            body = zlib.decompress(f.read())

            pos, out = 0, []
            while pos < len(body):
                (head,) = struct.unpack_from("<I", body, pos)
                pos += 4
                tag, size = head & 0x3FF, (head >> 20) & 0xFFF
                if size == 0xFFF:
                    (size,) = struct.unpack_from("<I", body, pos)
                    pos += 4
                data = body[pos : pos + size]
                pos += size
                if tag == 67:  # PARA_TEXT
                    try:
                        out.append(data.decode("utf-16le").rstrip("\x00"))
                    except UnicodeDecodeError:
                        pass
            txt = "\n".join(out)
        elif magic.startswith(b"<?xm") or magic.startswith(b"<DOC"):
            data = f.read().decode("utf-8", "ignore")
            root = ET.fromstring(data)
            out = [n.text for n in root.iter("CHAR") if n.text]
            txt = "\n".join(out)
        else:
            txt = ""

    return txt


# HWPX 문서에서 텍스트 추출
def extract_hwpx(path):
    texts = []
    with zipfile.ZipFile(path, "r") as z:
        for name in z.namelist():
            if name.startswith("Contents/section") and name.endswith(".xml"):
                content = z.read(name)
                root = ET.fromstring(content)
                for t_elem in root.iter():
                    if t_elem.tag.endswith("t") and t_elem.text:
                        try:
                            # 이스케이프된 한글 바이트인 경우 decode 시도
                            decoded = bytes(t_elem.text, "latin1").decode("utf-8")
                            texts.append(decoded)
                        except:
                            texts.append(t_elem.text.strip())
    return " ".join(texts)


# DOCX 문서에서 텍스트 추출
def extract_docx(path):
    try:
        doc = docx.Document(path)
        texts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
        return " ".join(texts)
    except Exception as e:
        return ""


# PDF 문서에서 텍스트 추출
def extract_pdf(path):
    try:
        doc = fitz.open(path)
        texts = []
        for page in doc:
            page_text = page.get_text("text")
            lines = [line.strip() for line in page_text.splitlines() if line.strip()]
            texts.extend(lines)
        return " ".join(texts)
    except Exception as e:
        return ""


def extract_pptx(filepath):
    prs = Presentation(filepath)
    all_texts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 일반 텍스트
                all_texts.append(shape.text.strip())
            elif shape.shape_type == 19:  # 표(table)
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            all_texts.append(cell.text.strip())

    # 모든 텍스트를 스페이스로 연결해서 한 줄로 저장
    flat_line = " ".join(all_texts)
    return flat_line


def extract_default(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        data = f.read()
        return data


# 파일 확장자에 따라 적절한 추출 함수 선택
def process_text(filepath):
    import jpype

    if not jpype.isJVMStarted():
        jar_dir = Path("/maven/target/dependency")
        jars = [str(p) for p in jar_dir.glob("*.jar")]
        jpype.startJVM(classpath=jars)
    import jpype.imports

    mime = magic.Magic(mime=True)
    mime_type = mime.from_file(filepath)
    ext = "default"
    if olefile.isOleFile(filepath):
        try:
            ole = olefile.OleFileIO(filepath)
            streams = ole.listdir()
            stream_names = [".".join(s) for s in streams]
            if any("PowerPoint Document" in s for s in stream_names):
                ext = "ppt"
            elif any("WordDocument" in s for s in stream_names):
                ext = "doc"
            elif any(
                s in stream_names for s in ["BodyText", "FileHeader", "HwpSummary"]
            ):
                ext = "hwp"
        except:
            pass
    elif mime_type == "application/zip":
        try:
            with zipfile.ZipFile(filepath, "r") as zipf:
                names = [name.replace("\\", "/") for name in zipf.namelist()]
                if any(name.endswith("word/document.xml") for name in names):
                    ext = "docx"
                elif any(name.endswith("ppt/presentation.xml") for name in names):
                    ext = "pptx"
                elif any(name.endswith("xl/workbook.xml") for name in names):
                    ext = "xlsx"
        except:
            pass
    elif (
        mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        try:
            with zipfile.ZipFile(filepath, "r") as zipf:
                names = [name.replace("\\", "/") for name in zipf.namelist()]
                if any(name.endswith("word/document.xml") for name in names):
                    ext = "docx"
        except:
            pass
    elif (
        mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        try:
            with zipfile.ZipFile(filepath, "r") as zipf:
                names = [name.replace("\\", "/") for name in zipf.namelist()]
                if any(name.endswith("ppt/presentation.xml") for name in names):
                    ext = "pptx"
        except:
            pass
    elif mime_type in [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ]:
        ext = "xls"
    elif mime_type == "text/xml":
        try:
            hwpfilepath = pathlib.Path(filepath).expanduser()
            with hwpfilepath.open("rb") as f:
                xml_bytes = f.read().decode("utf-8", "ignore")
                root = ET.fromstring(xml_bytes)
                tag = root.tag.split("}", 1)[-1]  # {namespace}Tag → Tag
                ns = root.tag.split("}", 1)[0][1:] if "}" in root.tag else ""
                if (
                    tag in {"HWPX", "HWPML", "HwpDocuments"}
                    or "hancom.co.kr/hwpml" in ns.lower()
                ):
                    ext = "hwp"
                else:
                    pass
        except ET.ParseError:
            pass
    try:
        with zipfile.ZipFile(filepath, "r") as zipf:
            names = zipf.namelist()
            if any(name.startswith("Contents/") for name in names):
                ext = "hwpx"
    except:
        pass

    with open(filepath, "rb") as f:
        header = f.read(5)
        if header == b"%PDF-":
            ext = "pdf"

    text_extractor = {
        "hwp": extract_hwp,
        "doc": extract_ppt_doc,
        "ppt": extract_ppt_doc,
        "hwpx": extract_hwpx,
        "docx": extract_docx,
        "pptx": extract_pptx,
        "pdf": extract_pdf,
        "xlsx": extract_xlsx,
        "xls": extract_xls,
        "default": extract_default,
    }.get(ext)

    meta_extractor = {
        "hwp": get_file_ctime_created_date,
        "doc": get_ole_created_date,
        "ppt": get_ole_created_date,
        "hwpx": get_hwpx_created_date,
        "docx": get_office_openxml_created_date,
        "pptx": get_office_openxml_created_date,
        "pdf": get_pdf_created_date,
        "xls": get_ole_created_date,
        "xlsx": get_office_openxml_created_date,
        "default": get_file_ctime_created_date,
    }.get(ext)

    if text_extractor and meta_extractor:
        return {
            "metadata": {"createdate": meta_extractor(filepath)},
            "text": text_extractor(filepath),
        }
    else:
        return {"metadata": {"createdate": ""}, "text": ""}
